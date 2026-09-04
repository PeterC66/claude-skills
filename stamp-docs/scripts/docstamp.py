#!/usr/bin/env python3
"""docstamp - keep a version + date stamp current in Markdown and PowerPoint documents.

The stamp is visible when the document is viewed or printed, and carries a short hash of
the document's own content *with the stamp removed*. That hash is what makes this safe to
run on every turn from a Stop hook: if the content hash is unchanged the file is not
rewritten at all, so mtimes do not churn and there is no risk of a hook loop.

  Markdown   two lines just after the H1:
                 <!-- docstamp v1.4 | 2026-07-27 | sha=3f9a1c2b -->
                 **v1.4** * updated 27 July 2026
             The comment is invisible in every renderer and is the machine anchor.

  PowerPoint a small text box on every slide, shape-named "docstamp" so a re-run updates
             it instead of stacking a second one. The stamp record lives in the file's
             core properties. Position adapts to whichever candidate band is clear on
             every slide; colour adapts to background luminance.

Usage:
  docstamp.py --auto            hook mode: mtime-gated scan, stamp what changed, always exit 0
  docstamp.py --all             force a full hash scan, ignoring the mtime gate
  docstamp.py --backfill        first run: stamp everything unstamped at v1.0, dated from git
  docstamp.py --check           audit only; exits 1 if anything is unstamped or stale
  docstamp.py --list            list the in-scope documents and exit
  docstamp.py --major FILE      bump the major version of one file (a rewrite)
  docstamp.py --minor FILE      force a minor bump of one file

  --dry-run        report what would change, write nothing
  --root NAME      restrict to one root from the policy (buses | portal | ops)
  --policy PATH    use an alternative policy file (used by the test harness)
"""

from __future__ import annotations

import argparse
import datetime as dt
import fnmatch
import hashlib
import json
import os
import re
import subprocess
import sys
import traceback

HERE = os.path.dirname(os.path.abspath(__file__))
SKILL_DIR = os.path.dirname(HERE)
DEFAULT_POLICY = os.path.join(SKILL_DIR, "stamp-policy.json")
LOG_PATH = os.path.join(SKILL_DIR, "docstamp.log")

MIDDOT = "·"  # the separator used in the visible stamp
STAMP_SHAPE_NAME = "docstamp"

MD_COMMENT_RE = re.compile(
    r"^<!--\s*docstamp\s+v(\d+)\.(\d+)\s*\|\s*(\d{4}-\d{2}-\d{2})\s*\|\s*sha=([0-9a-fA-F]+)\s*-->\s*$"
)
MD_VISIBLE_RE = re.compile(r"^\*\*v\d+\.\d+\*\*\s*.\s*updated\s+\S.*$")
PPTX_PROP_RE = re.compile(
    r"docstamp\s+v(\d+)\.(\d+)\s*\|\s*(\d{4}-\d{2}-\d{2})\s*\|\s*sha=([0-9a-fA-F]+)"
)


# --------------------------------------------------------------------------- helpers


def human_date(d: dt.date) -> str:
    """27 July 2026 - written out rather than %-d/%#d, which is not portable."""
    return "{} {} {}".format(d.day, d.strftime("%B"), d.year)


def stamp_comment(major, minor, iso, sha):
    return "<!-- docstamp v{}.{} | {} | sha={} -->".format(major, minor, iso, sha)


def stamp_visible(major, minor, d):
    return "**v{}.{}** {} updated {}".format(major, minor, MIDDOT, human_date(d))


def pptx_stamp_text(major, minor, d):
    return "v{}.{} {} {}".format(major, minor, MIDDOT, human_date(d))


def load_policy(path):
    with open(path, "r", encoding="utf-8") as fh:
        return json.load(fh)


def state_path(policy_file):
    return os.path.join(os.path.dirname(policy_file), ".docstamp-state.json")


def load_state(policy_file):
    try:
        with open(state_path(policy_file), "r", encoding="utf-8") as fh:
            return json.load(fh)
    except Exception:
        return {}


def save_state(policy_file, state):
    with open(state_path(policy_file), "w", encoding="utf-8") as fh:
        json.dump(state, fh, indent=2)


def git_date(root, abs_path):
    """Date of the last commit touching this file, or None if untracked / not a repo."""
    try:
        out = subprocess.run(
            ["git", "-C", root, "log", "-1", "--format=%cs", "--", abs_path],
            capture_output=True, text=True, timeout=20,
        )
        val = out.stdout.strip()
        if re.fullmatch(r"\d{4}-\d{2}-\d{2}", val):
            return dt.date.fromisoformat(val)
    except Exception:
        pass
    return None


def mtime_date(abs_path):
    return dt.date.fromtimestamp(os.path.getmtime(abs_path))


# --------------------------------------------------------------------------- discovery


# A document written by a SCRIPT rather than by a person, recognised by a marker
# in its own head rather than by its path.
#
# WHY A CONTENT RULE AND NOT ANOTHER GLOB (open action OA-096). A generated report
# and a stamp fight each other: the Stop hook inserts a docstamp block, re-running
# the generator overwrites the file and deletes it, the hook puts it back, and the
# file churns in every diff while its "version" describes nothing. Two such reports
# were doing exactly that -- `stand-coverage_2026-08-22.md` and, since 2026-08-23,
# `frame-coverage_2026-08-23.md`.
#
# Naming those two in `excludeGlobs` would have fixed those two. The THIRD generated
# report would arrive with the same fault, because nothing about writing a generator
# reminds anybody that a policy file in another repository has to learn its filename
# -- which is how both of these got here in the first place. Keying on a marker the
# report itself carries makes the rule self-applying: a generator that says it
# generates the file is excluded by having said so, on the day it is written, by
# somebody who is already thinking about that file.
#
# The marker was not invented for this. `frame-coverage_2026-08-23.md` already
# carried `<!-- generated by measure_frame_coverage_2026-08-23.py v1.0; ... -->`,
# written for a human reader; this makes it load-bearing as well.
#
# Read from the HEAD of the file only, so that a document merely discussing
# generated reports -- this project has several -- does not exclude itself by
# quoting one. 40 lines is comfortably past any front matter and title.
GENERATED_MARKER_LINES = 40


def is_generated(abs_path, marker):
    if not marker or not abs_path.lower().endswith(".md"):
        return False
    try:
        with open(abs_path, "r", encoding="utf-8", errors="replace") as fh:
            for i, line in enumerate(fh):
                if i >= GENERATED_MARKER_LINES:
                    return False
                if marker in line:
                    return True
    except OSError:
        return False
    return False


def discover(policy, only_root=None):
    """Walk the configured roots, pruning excluded directories before descending."""
    found = []
    for root_cfg in policy["roots"]:
        if only_root and root_cfg["name"] != only_root:
            continue
        root = os.path.normpath(root_cfg["path"])
        if not os.path.isdir(root):
            continue
        exts = tuple(e.lower() for e in root_cfg["extensions"])
        skip_names = set(root_cfg.get("excludeDirNames", []))
        skip_pats = root_cfg.get("excludeDirPatterns", [])
        skip_globs = root_cfg.get("excludeGlobs", [])
        generated_marker = policy.get("generatedMarker")

        for dirpath, dirnames, filenames in os.walk(root):
            dirnames[:] = [
                d for d in dirnames
                if d not in skip_names
                and not any(fnmatch.fnmatch(d, p) for p in skip_pats)
            ]
            for fn in filenames:
                if not fn.lower().endswith(exts):
                    continue
                if fn.startswith("~$"):  # Office lock files
                    continue
                abs_path = os.path.join(dirpath, fn)
                rel = os.path.relpath(abs_path, root).replace(os.sep, "/")
                if any(fnmatch.fnmatch(rel, g) for g in skip_globs):
                    continue
                if is_generated(abs_path, generated_marker):
                    continue
                found.append((root_cfg["name"], root, abs_path, rel))
    found.sort(key=lambda t: (t[0], t[3].lower()))
    return found


# --------------------------------------------------------------------------- markdown


def read_text(path):
    with open(path, "rb") as fh:
        raw = fh.read()
    bom = raw.startswith(b"\xef\xbb\xbf")
    if bom:
        raw = raw[3:]
    newline = "\r\n" if b"\r\n" in raw else "\n"
    return raw.decode("utf-8").replace("\r\n", "\n"), newline, bom


def write_text(path, text, newline, bom):
    data = text.replace("\n", newline).encode("utf-8")
    if bom:
        data = b"\xef\xbb\xbf" + data
    with open(path, "wb") as fh:
        fh.write(data)


def md_strip_stamp(lines):
    """Remove an existing stamp block. Returns (lines_without, existing_meta_or_None)."""
    for i, line in enumerate(lines):
        m = MD_COMMENT_RE.match(line)
        if not m:
            continue
        meta = {
            "major": int(m.group(1)), "minor": int(m.group(2)),
            "date": m.group(3), "sha": m.group(4).lower(),
        }
        j = i + 1
        if j < len(lines) and MD_VISIBLE_RE.match(lines[j]):
            j += 1
        out = lines[:i] + lines[j:]
        # collapse the double blank line left behind at the seam
        if 0 < i < len(out) and out[i - 1].strip() == "" and out[i].strip() == "":
            del out[i]
        return out, meta
    return list(lines), None


def md_body_hash(lines, sha_len):
    body = "\n".join(ln.rstrip() for ln in lines).strip() + "\n"
    return hashlib.sha256(body.encode("utf-8")).hexdigest()[:sha_len]


def md_insert_at(lines):
    """Index to insert the stamp at: after YAML frontmatter and the H1, if present."""
    start = 0
    if lines and lines[0].strip() == "---":
        for k in range(1, min(len(lines), 60)):
            if lines[k].strip() == "---":
                start = k + 1
                break
    for k in range(start, min(len(lines), start + 20)):
        if re.match(r"^#\s+\S", lines[k]):
            return k + 1
    return start


def stamp_markdown(abs_path, sha_len, when, force=None, dry_run=False):
    text, newline, bom = read_text(abs_path)
    lines = text.split("\n")
    stripped, existing = md_strip_stamp(lines)
    sha = md_body_hash(stripped, sha_len)

    if existing and existing["sha"] == sha and not force:
        return ("unchanged", existing["major"], existing["minor"])

    if existing is None:
        major, minor = 1, 0
    elif force == "major":
        major, minor = existing["major"] + 1, 0
    else:
        major, minor = existing["major"], existing["minor"] + 1

    at = md_insert_at(stripped)
    tail = stripped[at:]
    while tail and tail[0].strip() == "":
        tail.pop(0)
    block = [stamp_comment(major, minor, when.isoformat(), sha), stamp_visible(major, minor, when)]
    head = stripped[:at]
    if head:
        new_lines = head + [""] + block + [""] + tail
    else:
        new_lines = block + [""] + tail

    if not dry_run:
        write_text(abs_path, "\n".join(new_lines), newline, bom)
    return ("stamped" if existing is None else "bumped", major, minor)


# ------------------------------------------------------------------------- powerpoint


def _iter_shapes(shapes):
    for sp in shapes:
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            if sp.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from _iter_shapes(sp.shapes)
                continue
        except Exception:
            pass
        yield sp


def pptx_digest(prs, sha_len):
    """Slide text in order, plus every embedded media blob, so a swapped image counts."""
    h = hashlib.sha256()
    for si, slide in enumerate(prs.slides):
        h.update("\x01slide{}".format(si).encode("utf-8"))
        for sp in _iter_shapes(slide.shapes):
            if sp.name == STAMP_SHAPE_NAME:
                continue
            try:
                if sp.has_text_frame:
                    h.update(sp.text_frame.text.encode("utf-8"))
            except Exception:
                pass
            try:
                if sp.has_table:
                    for row in sp.table.rows:
                        for cell in row.cells:
                            h.update(cell.text.encode("utf-8"))
            except Exception:
                pass
    try:
        parts = sorted(prs.part.package.iter_parts(), key=lambda p: str(p.partname))
        for part in parts:
            name = str(part.partname)
            if name.startswith("/ppt/media/"):
                h.update(name.encode("utf-8"))
                h.update(hashlib.sha256(part.blob).digest())
    except Exception:
        pass
    return h.hexdigest()[:sha_len]


def candidate_rects(width, height):
    emu = 914400
    w, h = int(3.2 * emu), int(0.24 * emu)
    margin = int(0.30 * emu)
    return {
        "top-right":     (width - w - margin, int(0.10 * emu),          w, h, "right"),
        "bottom-right":  (width - w - margin, height - h - int(0.14 * emu), w, h, "right"),
        "bottom-left":   (margin,             height - h - int(0.14 * emu), w, h, "left"),
        "bottom-centre": (int((width - w) / 2), height - h - int(0.14 * emu), w, h, "center"),
    }


def _clashes(slide, rect):
    rl, rt, rw, rh = rect[:4]
    hits = []
    for sp in slide.shapes:
        if sp.name == STAMP_SHAPE_NAME:
            continue
        l, t, w, h = sp.left, sp.top, sp.width, sp.height
        if None in (l, t, w, h):
            continue
        if l < rl + rw and l + w > rl and t < rt + rh and t + h > rt:
            hits.append(sp.name)
    return hits


def choose_position(prs, preference):
    rects = candidate_rects(prs.slide_width, prs.slide_height)
    for name in preference:
        rect = rects.get(name)
        if rect and not any(_clashes(s, rect) for s in prs.slides):
            return name, rect, []
    name = preference[0]
    rect = rects[name]
    clashing = [i for i, s in enumerate(prs.slides, 1) if _clashes(s, rect)]
    return name, rect, clashing


def _resolve_bg_rgb(slide):
    from pptx.enum.dml import MSO_FILL
    for src in (slide, slide.slide_layout, slide.slide_layout.slide_master):
        try:
            fill = src.background.fill
            if fill.type == MSO_FILL.SOLID:
                return fill.fore_color.rgb
        except Exception:
            continue
    return None


def _stamp_colour(slide, cfg):
    from pptx.dml.color import RGBColor
    rgb = _resolve_bg_rgb(slide)
    if rgb is None:
        return RGBColor.from_string(cfg.get("colourOnLight", "7A7A7A"))
    r, g, b = rgb[0], rgb[1], rgb[2]
    lum = 0.2126 * r + 0.7152 * g + 0.0722 * b
    key = "colourOnDark" if lum < 128 else "colourOnLight"
    return RGBColor.from_string(cfg[key])


def _remove_existing_stamps(slide):
    for sp in list(slide.shapes):
        if sp.name == STAMP_SHAPE_NAME:
            sp._element.getparent().remove(sp._element)


def stamp_pptx(abs_path, sha_len, when, cfg, force=None, dry_run=False):
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    prs = Presentation(abs_path)
    sha = pptx_digest(prs, sha_len)

    existing = None
    try:
        m = PPTX_PROP_RE.search(prs.core_properties.comments or "")
        if m:
            existing = {
                "major": int(m.group(1)), "minor": int(m.group(2)),
                "date": m.group(3), "sha": m.group(4).lower(),
            }
    except Exception:
        pass

    if existing and existing["sha"] == sha and not force:
        return ("unchanged", existing["major"], existing["minor"], None, [])

    if existing is None:
        major, minor = 1, 0
    elif force == "major":
        major, minor = existing["major"] + 1, 0
    else:
        major, minor = existing["major"], existing["minor"] + 1

    pos_name, rect, clashing = choose_position(prs, cfg["positionPreference"])
    left, top, width, height, align = rect
    text = pptx_stamp_text(major, minor, when)
    alignment = {"right": PP_ALIGN.RIGHT, "left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER}[align]

    if not dry_run:
        for slide in prs.slides:
            _remove_existing_stamps(slide)
            box = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
            box.name = STAMP_SHAPE_NAME
            tf = box.text_frame
            tf.word_wrap = False
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            para = tf.paragraphs[0]
            para.alignment = alignment
            run = para.add_run()
            run.text = text
            run.font.size = Pt(cfg.get("fontSizePt", 9))
            run.font.bold = False
            run.font.color.rgb = _stamp_colour(slide, cfg)
        prs.core_properties.comments = "docstamp v{}.{} | {} | sha={}".format(
            major, minor, when.isoformat(), sha
        )
        prs.save(abs_path)

    return ("stamped" if existing is None else "bumped", major, minor, pos_name, clashing)


# ------------------------------------------------------------------------------- run


def check_one(abs_path, sha_len, cfg):
    """Audit a single file. Returns (state, version_str) without writing anything."""
    if abs_path.lower().endswith(".md"):
        text, _, _ = read_text(abs_path)
        stripped, existing = md_strip_stamp(text.split("\n"))
        if existing is None:
            return "missing", None
        sha = md_body_hash(stripped, sha_len)
        ver = "v{}.{}".format(existing["major"], existing["minor"])
        return ("current" if existing["sha"] == sha else "stale"), ver
    from pptx import Presentation
    prs = Presentation(abs_path)
    m = PPTX_PROP_RE.search(prs.core_properties.comments or "")
    if not m:
        return "missing", None
    sha = pptx_digest(prs, sha_len)
    ver = "v{}.{}".format(m.group(1), m.group(2))
    return ("current" if m.group(4).lower() == sha else "stale"), ver


def process(args, policy):
    sha_len = policy.get("shaLen", 8)
    pcfg = policy.get("pptx", {})
    files = discover(policy, args.root)
    today = dt.date.today()

    # --auto only hashes files that look newer than the last successful run.
    gate = None
    if args.auto and not args.all:
        st = load_state(args.policy)
        last = st.get("lastRun")
        if last:
            gate = last - 5  # small clock-skew allowance

    results = []
    for name, root, abs_path, rel in files:
        if gate is not None:
            try:
                if os.path.getmtime(abs_path) <= gate:
                    continue
            except OSError:
                continue

        if args.check:
            state, ver = check_one(abs_path, sha_len, pcfg)
            results.append((state, name, rel, ver, None, []))
            continue

        when = today
        if args.backfill and check_one(abs_path, sha_len, pcfg)[0] == "missing":
            # Backfill dates a NEW document from git, because "updated" should mean
            # when it was actually last written, not when backfill happened to run.
            # It must not do that to a document that already carries a stamp: an
            # already-stamped file only reaches the stamping code because its content
            # CHANGED, so its git date is the previous commit and dating it from there
            # walks the visible "updated" line backwards. On 2026-08-21 a --backfill run
            # for one new README re-dated four edited ones to the day before.
            when = git_date(root, abs_path) or mtime_date(abs_path)

        if abs_path.lower().endswith(".md"):
            state, major, minor = stamp_markdown(
                abs_path, sha_len, when, force=args.force, dry_run=args.dry_run
            )
            pos, clash = None, []
        else:
            state, major, minor, pos, clash = stamp_pptx(
                abs_path, sha_len, when, pcfg, force=args.force, dry_run=args.dry_run
            )
        results.append((state, name, rel, "v{}.{}".format(major, minor), pos, clash))

    return results


def main(argv=None):
    ap = argparse.ArgumentParser(add_help=True)
    mode = ap.add_mutually_exclusive_group()
    mode.add_argument("--auto", action="store_true", help="hook mode; always exits 0")
    mode.add_argument("--check", action="store_true", help="audit only; exits 1 on findings")
    mode.add_argument("--backfill", action="store_true", help="stamp unstamped docs, dated from git")
    mode.add_argument("--list", action="store_true", help="list in-scope documents")
    ap.add_argument("--all", action="store_true", help="ignore the mtime gate; full hash scan")
    ap.add_argument("--major", metavar="FILE", help="bump the major version of one file")
    ap.add_argument("--minor", metavar="FILE", help="force a minor bump of one file")
    ap.add_argument("--dry-run", action="store_true")
    ap.add_argument("--root", help="restrict to one root name from the policy")
    ap.add_argument("--policy", default=DEFAULT_POLICY)
    ap.add_argument("--quiet", action="store_true")
    args = ap.parse_args(argv)
    args.force = None

    policy = load_policy(args.policy)
    sha_len = policy.get("shaLen", 8)

    if args.list:
        for name, _root, _abs, rel in discover(policy, args.root):
            print("{:8s} {}".format(name, rel))
        return 0

    # single-file forced bump
    target = args.major or args.minor
    if target:
        args.force = "major" if args.major else "minor"
        abs_path = os.path.abspath(target)
        today = dt.date.today()
        if abs_path.lower().endswith(".md"):
            state, maj, mn = stamp_markdown(abs_path, sha_len, today, force=args.force,
                                            dry_run=args.dry_run)
        else:
            state, maj, mn, _pos, _cl = stamp_pptx(abs_path, sha_len, today,
                                                   policy.get("pptx", {}),
                                                   force=args.force, dry_run=args.dry_run)
        print("{}  v{}.{}  {}".format(state, maj, mn, os.path.basename(abs_path)))
        return 0

    run_started = dt.datetime.now().timestamp()
    results = process(args, policy)

    if args.check:
        bad = [r for r in results if r[0] in ("missing", "stale")]
        if not args.quiet:
            for state, root, rel, ver, _pos, _cl in results:
                if state != "current":
                    print("{:8s} {:8s} {}".format(state, root, rel))
            print("docstamp --check: {} document(s), {} need attention".format(len(results), len(bad)))
        return 1 if bad else 0

    changed = [r for r in results if r[0] in ("stamped", "bumped")]
    clashes = [(r[2], r[5]) for r in results if r[5]]

    if args.auto:
        if changed and not args.quiet:
            head = ", ".join("{} {}".format(os.path.basename(r[2]), r[3]) for r in changed[:3])
            more = "" if len(changed) <= 3 else " (+{} more)".format(len(changed) - 3)
            print("docstamp: stamped {} document(s) - {}{}".format(len(changed), head, more))
        if not args.dry_run:
            save_state(args.policy, {"lastRun": run_started})
        return 0

    if not args.quiet:
        for state, root, rel, ver, pos, clash in results:
            if state in ("stamped", "bumped"):
                extra = "  [{}]".format(pos) if pos else ""
                if clash:
                    extra += "  CLASH on slides {}".format(clash)
                print("{:9s} {:6s} {:8s} {}{}".format(state, ver, root, rel, extra))
        print("---- {} document(s): {} changed, {} already current".format(
            len(results), len(changed), len(results) - len(changed)))
        for rel, clash in clashes:
            print("WARNING: no clear stamp position in {} (slides {})".format(rel, clash))
    if not args.dry_run and not args.check:
        save_state(args.policy, {"lastRun": run_started})
    return 0


if __name__ == "__main__":
    hook_mode = "--auto" in sys.argv
    try:
        sys.exit(main())
    except SystemExit:
        raise
    except Exception:
        # A hook must never fail the turn: log and exit clean.
        try:
            with open(LOG_PATH, "a", encoding="utf-8") as fh:
                fh.write("\n=== {} ===\n{}".format(dt.datetime.now().isoformat(),
                                                   traceback.format_exc()))
        except Exception:
            pass
        if hook_mode:
            sys.exit(0)
        traceback.print_exc()
        sys.exit(2)
